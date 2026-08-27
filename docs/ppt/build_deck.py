#!/usr/bin/env python3
"""
바이브키 발표자료 생성기 (2026 임베디드SW 경진대회 자유공모)

  python3 docs/ppt/build_deck.py

수치를 손으로 옮겨 적지 않습니다. 아래 FACTS 한 곳만 고치면 전체 장표가 따라갑니다.
지어낸 값을 넣지 않기 위해, 아직 측정하지 않은 항목은 MEASURED = None 으로 두고
장표에 "측정 예정"으로 표시합니다.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------------------------------------------------------------- 사실 (한 곳에서 관리)
FACTS = {
    "team": "뚝딱뚝딱",
    "product": "바이브키",
    "contest": "2026 임베디드SW 경진대회 · 자유공모 부문",
    "mcu": "Seeed Studio XIAO ESP32-S3",
    "tests_fw": 64,
    "tests_app": 48,
    "flash_used": 287025,
    "flash_total": 3342336,
    "ram_used": 22888,
    "ram_total": 327680,
    "bom_won": 14690,
    "bom_items": 5,
    "buttons": 3,
}
FACTS["tests_total"] = FACTS["tests_fw"] + FACTS["tests_app"]
FACTS["flash_pct"] = FACTS["flash_used"] / FACTS["flash_total"] * 100
FACTS["ram_pct"] = FACTS["ram_used"] / FACTS["ram_total"] * 100

# 실기기 측정값. 측정 세션 후 채웁니다. None 이면 장표에 "측정 예정"으로 나갑니다.
MEASURED = {
    "latency_avg_ms": None,
    "latency_max_ms": None,
    "frame_error_pct": None,
    "misfire_count": None,
    "idle_current_ma": None,
    "sleep_current_ua": None,
}

# ---------------------------------------------------------------- 색·서체
BG     = RGBColor(0xF7, 0xF9, 0xFC)
INK    = RGBColor(0x10, 0x14, 0x18)
COBALT = RGBColor(0x0B, 0x57, 0xD0)
MUTED  = RGBColor(0x41, 0x50, 0x5E)
LIGHT  = RGBColor(0x8B, 0x93, 0xA1)
GREEN  = RGBColor(0x12, 0x70, 0x3A)
AMBER  = RGBColor(0xB2, 0x6A, 0x00)
RED    = RGBColor(0xB3, 0x26, 0x1E)
CARD   = RGBColor(0xED, 0xF2, 0xFB)
BORDER = RGBColor(0xC9, 0xD9, 0xF2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

SANS = "Apple SD Gothic Neo"   # 윈도우에서 발표한다면 "맑은 고딕" 으로 바꾸세요
MONO = "Menlo"

W, H = Inches(13.333), Inches(7.5)


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def slide(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])   # 빈 레이아웃
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK if dark else BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    """runs = [(문자열, 크기pt, 굵기, 색, 폰트)] — 한 문단 안의 조각들"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for line in runs:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = align
        para.line_spacing = spacing
        for (txt, size, bold, color, font) in line:
            r = para.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return tb


def one(s, x, y, w, h, txt, size, bold=False, color=INK, font=SANS,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.15):
    return text(s, x, y, w, h, [[(txt, size, bold, color, font)]], align, anchor, spacing)


def card(s, x, y, w, h, fill=CARD, line=BORDER):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1.25)
    box.shadow.inherit = False
    box.adjustments[0] = 0.06
    box.text_frame.text = ""
    return box


def header(s, eyebrow, title, dark=False):
    one(s, Inches(0.9), Inches(0.62), Inches(11.5), Inches(0.3),
        eyebrow, 13, True, COBALT if not dark else RGBColor(0x6F, 0xA8, 0xF5))
    one(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.8),
        title, 34, True, INK if not dark else BG)


def rule(s, y, x=Inches(0.9), w=Inches(11.5), color=BORDER):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(1.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background(); ln.shadow.inherit = False


def fmt(v, unit, digits=2):
    return "측정 예정" if v is None else f"{v:.{digits}f} {unit}".replace(".00 ", " ")


# ================================================================ 장표
prs = deck()

# --- 1. 표지 -----------------------------------------------------
s = slide(prs, dark=True)
one(s, Inches(1.1), Inches(2.35), Inches(11), Inches(0.4),
    FACTS["contest"], 16, True, RGBColor(0x6F, 0xA8, 0xF5))
one(s, Inches(1.1), Inches(2.85), Inches(11), Inches(1.75),
    FACTS["product"], 78, True, BG)
one(s, Inches(1.1), Inches(4.25), Inches(11), Inches(0.6),
    "어르신도 버튼 세 번이면 끝", 28, False, RGBColor(0xC7, 0xD3, 0xE0))
rule(s, Inches(5.1), Inches(1.1), Inches(2.2), RGBColor(0x6F, 0xA8, 0xF5))
one(s, Inches(1.1), Inches(5.45), Inches(11), Inches(0.4),
    f"팀 {FACTS['team']}", 18, True, BG)
one(s, Inches(1.1), Inches(5.85), Inches(11), Inches(0.4),
    "USB 물리 버튼 입력기 · 펌웨어 3.1 / 앱 3.1", 14, False, LIGHT, MONO)

# --- 2. 문제 -----------------------------------------------------
s = slide(prs)
header(s, "왜 만들었나", "스마트폰이 어려운 분에게 화면은 벽입니다")
rows = [
    ("아이콘 스무 개 중에서 찾아야 합니다", "이름이 작고, 비슷하게 생겼습니다"),
    ("화면을 세 번 넘겨야 전화가 걸립니다", "중간에 한 번만 잘못 눌러도 처음부터"),
    ("잘못 누르면 무슨 일이 생길지 모릅니다", "그래서 아예 시도하지 않게 됩니다"),
]
y = Inches(2.25)
for head, sub in rows:
    c = card(s, Inches(0.9), y, Inches(11.5), Inches(1.05))
    one(s, Inches(1.35), y + Inches(0.2), Inches(10.6), Inches(0.4), head, 21, True, INK)
    one(s, Inches(1.35), y + Inches(0.62), Inches(10.6), Inches(0.35), sub, 15, False, MUTED)
    y += Inches(1.25)
one(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.5),
    "→ 문제는 앱이 부족한 것이 아니라, 입력 수단이 화면 하나뿐이라는 점입니다.",
    19, True, COBALT)

# --- 3. 해답 -----------------------------------------------------
s = slide(prs)
header(s, "무엇을 만들었나", "버튼 세 개, 케이블 하나")
cols = [("1", "전화", COBALT), ("2", "문자", GREEN), ("3", "길찾기", AMBER)]
x = Inches(0.9)
for num, label, col in cols:
    card(s, x, Inches(2.2), Inches(3.55), Inches(2.0))
    one(s, x, Inches(2.5), Inches(3.55), Inches(1.15), num, 54, True, col, MONO, PP_ALIGN.CENTER)
    one(s, x, Inches(3.42), Inches(3.55), Inches(0.5), label, 24, True, INK, SANS, PP_ALIGN.CENTER)
    x += Inches(3.98)
one(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(0.5),
    "USB로 꽂으면 끝입니다. 인터넷도, 계정도, 충전도 필요 없습니다.", 22, True, INK)
facts = [
    ("MCU", FACTS["mcu"]),
    ("연결", "USB CDC (네이티브 USB, 0x303A:0x1001)"),
    ("전원", "폰에서 공급 — 배터리 없음"),
    ("부품", f"{FACTS['bom_items']}종 · 약 {FACTS['bom_won']:,}원"),
]
y = Inches(5.25)
for k, v in facts:
    one(s, Inches(0.95), y, Inches(1.6), Inches(0.3), k, 14, True, COBALT)
    one(s, Inches(2.6), y, Inches(9.6), Inches(0.3), v, 14, False, MUTED)
    y += Inches(0.42)

# --- 4. 쓰는 법 --------------------------------------------------
s = slide(prs)
header(s, "어떻게 쓰나", "버튼 3개로 6가지 — 짧게와 길게를 구분합니다")
card(s, Inches(0.9), Inches(2.2), Inches(5.55), Inches(3.1))
one(s, Inches(1.3), Inches(2.5), Inches(4.8), Inches(0.4), "짧게 누르기", 20, True, COBALT)
one(s, Inches(1.3), Inches(3.0), Inches(4.8), Inches(1.9),
    "1번 → 전화 앱 (번호까지 채워서)\n2번 → 문자 앱\n3번 → 길찾기", 18, False, INK, SANS, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 1.6)
card(s, Inches(6.85), Inches(2.2), Inches(5.55), Inches(3.1))
one(s, Inches(7.25), Inches(2.5), Inches(4.8), Inches(0.4), "0.7초 이상 길게", 20, True, AMBER)
one(s, Inches(7.25), Inches(3.0), Inches(4.8), Inches(1.9),
    "1번 → 사용자가 지정한 앱\n2번 → 사용자가 지정한 앱\n3번 → AI 도우미 (음성)", 18, False, INK, SANS, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 1.6)
one(s, Inches(0.9), Inches(5.65), Inches(11.5), Inches(0.6),
    "무엇을 넣을지는 쓰는 사람이 정합니다. 앱에서 끌어다 놓으면 됩니다.", 19, True, MUTED)
one(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.5),
    "버튼을 늘리지 않고 할 수 있는 일을 늘렸습니다 — 늘어난 버튼은 곧 늘어난 혼란이기 때문입니다.",
    16, False, LIGHT)

# --- 5. 시스템 구성 ----------------------------------------------
s = slide(prs)
header(s, "시스템 구성", "세 층이 각자 책임을 집니다")
layers = [
    ("기기 (ESP32-S3)", "버튼을 인터럽트로 잡고, 채터링을 거르고, 프레임으로 포장", COBALT),
    ("VKP v1 프로토콜", "CRC16 서명 · ACK 재전송 · SEQ 중복 제거", GREEN),
    ("안드로이드 앱", "프레임 해석 · 앱 실행 · 진동 피드백 · 자가진단", AMBER),
]
y = Inches(2.25)
for name, desc, col in layers:
    card(s, Inches(0.9), y, Inches(11.5), Inches(1.25))
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), y, Pt(6), Inches(1.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = col
    bar.line.fill.background(); bar.shadow.inherit = False
    one(s, Inches(1.35), y + Inches(0.24), Inches(10.6), Inches(0.4), name, 22, True, INK)
    one(s, Inches(1.35), y + Inches(0.72), Inches(10.6), Inches(0.35), desc, 15, False, MUTED)
    y += Inches(1.45)
one(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
    "기기는 '전선'이 아니라 스스로 판단하고 책임지는 노드입니다.", 18, True, COBALT)

# --- 6. 기기 내부 ------------------------------------------------
s = slide(prs)
header(s, "기기 내부", "폴링하지 않습니다 — 눌린 순간을 인터럽트로 잡습니다")
stations = ["버튼", "인터럽트", "큐", "입력 태스크", "프로토콜 태스크", "USB"]
subs = ["D0·D1·D2", "IRAM_ATTR", "FreeRTOS", "10ms · core 1", "5ms · core 1", "CDC"]
x = Inches(0.72)
wcard = Inches(1.86)
for i, (st, sb) in enumerate(zip(stations, subs)):
    card(s, x, Inches(2.5), wcard, Inches(1.35))
    one(s, x, Inches(2.78), wcard, Inches(0.45), st, 16, True, INK, SANS, PP_ALIGN.CENTER)
    one(s, x, Inches(3.28), wcard, Inches(0.35), sb, 11, False, MUTED, MONO, PP_ALIGN.CENTER)
    if i < len(stations) - 1:
        one(s, x + wcard, Inches(2.95), Inches(0.28), Inches(0.4), "→", 18, True, COBALT, SANS, PP_ALIGN.CENTER)
    x += wcard + Inches(0.28)
pts = [
    "인터럽트가 눌린 시각을 마이크로초 단위로 기록합니다 — 지연을 잴 수 있는 이유입니다.",
    "채터링(접점 떨림) 25ms는 입력 태스크가 걸러 냅니다. 떨리는 손도 한 번으로 셉니다.",
    "태스크 두 개를 코어 1에 고정해, 입력 처리와 통신이 서로를 기다리지 않습니다.",
]
y = Inches(4.35)
for p in pts:
    one(s, Inches(0.95), y, Inches(0.3), Inches(0.35), "·", 18, True, COBALT)
    one(s, Inches(1.25), y, Inches(11.1), Inches(0.5), p, 16, False, MUTED)
    y += Inches(0.62)

# --- 7. 프로토콜 -------------------------------------------------
s = slide(prs)
header(s, "직접 만든 통신 규격", "VKP v1 — 전선으로 그냥 보내지 않습니다")
blocks = [("AA", "시작", 1.25), ("SEQ", "순번", 1.35), ("TYPE", "종류", 1.45),
          ("LEN", "길이", 1.25), ("DATA", "버튼·방식", 2.4), ("CRC16", "검사값", 1.8), ("55", "끝", 1.25)]
x = Inches(0.9)
for hex_, name, w in blocks:
    ww = Inches(w)
    card(s, x, Inches(2.45), ww, Inches(1.5))
    one(s, x, Inches(2.72), ww, Inches(0.55), hex_, 22, True, COBALT, MONO, PP_ALIGN.CENTER)
    one(s, x, Inches(3.32), ww, Inches(0.35), name, 13, False, MUTED, SANS, PP_ALIGN.CENTER)
    x += ww + Inches(0.13)
ul = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.28), Inches(4.12), Inches(6.6), Pt(4))
ul.fill.solid(); ul.fill.fore_color.rgb = COBALT
ul.line.fill.background(); ul.shadow.inherit = False
one(s, Inches(2.28), Inches(4.3), Inches(6.6), Inches(0.35),
    "이 구간을 CRC-16/CCITT-FALSE 로 계산해 서명합니다", 14, True, COBALT)
rule(s, Inches(5.0))
specs = [
    ("프레임 길이", "7 + LEN 바이트 (헤더 4 · CRC 2 · 끝 1)"),
    ("검사", "CRC-16/CCITT-FALSE (다항식 0x1021, 초기값 0xFFFF)"),
    ("신뢰", "ACK 미수신 시 최대 3회 재전송 · SEQ 동일 → 폰이 중복 실행 안 함"),
    ("호환", "구버전 평문 펌웨어도 그대로 동작 (앱이 미분류 바이트로 처리)"),
]
y = Inches(5.25)
for k, v in specs:
    one(s, Inches(0.95), y, Inches(2.0), Inches(0.3), k, 14, True, COBALT)
    one(s, Inches(3.0), y, Inches(9.3), Inches(0.3), v, 14, False, MUTED)
    y += Inches(0.42)

# --- 8. 신뢰성 ---------------------------------------------------
s = slide(prs)
header(s, "신뢰성", "틀어진 신호가 앱에 닿지 않게 합니다")
one(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.4),
    "손상 주입 시험 — 프레임에 무작위로 비트를 뒤집고 10만 번 반복", 17, True, INK)
tbl_rows = [
    ("손상 정도", "CRC8 (구버전)", "CRC16 (현재)"),
    ("1~3 비트", "0건  (0.000 %)", "0건  (0.000 %)"),
    ("4~10 비트", "8건  (0.008 %)", "0건  (0.000 %)"),
]
y = Inches(2.55)
for i, row in enumerate(tbl_rows):
    is_head = (i == 0)
    if not is_head:
        card(s, Inches(0.9), y, Inches(11.5), Inches(0.62), WHITE if i % 2 else CARD)
    xs = [Inches(1.3), Inches(5.3), Inches(8.9)]
    for j, cell in enumerate(row):
        col = COBALT if is_head else (INK if j == 0 else (RED if ("8건" in cell) else GREEN))
        one(s, xs[j], y + (Inches(0.0) if is_head else Inches(0.16)), Inches(3.6), Inches(0.4),
            cell, 15 if is_head else 16, True, col, SANS if j == 0 or is_head else MONO)
    y += Inches(0.5 if is_head else 0.72)
one(s, Inches(0.9), Inches(4.62), Inches(11.5), Inches(0.9),
    "CRC8 이 약해서 바꾼 것이 아닙니다. 1~3 비트 손상은 둘 다 완벽하게 잡습니다.\n"
    "여러 비트가 한꺼번에 틀어지는 경우에만 차이가 났고, 그 차이가 오실행이므로 CRC16 을 택했습니다.",
    15, False, MUTED, SANS, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 1.5)
one(s, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.4),
    "근거: firmware/test/test_vkp.cpp — 고정 시드로 언제든 재현됩니다", 14, False, LIGHT, MONO)

# --- 9. 자기 감시 ------------------------------------------------
s = slide(prs)
header(s, "스스로를 지켜본다", "멈추면 되살아나고, 왜 멈췄는지 보고합니다")
items = [
    ("태스크 워치독 5초", "태스크가 5초간 응답하지 않으면 자동 재시작", COBALT),
    ("재시작 원인 보고", "전원·워치독·패닉·USB 등 16가지를 구분해 폰에 전달", GREEN),
    ("놓친 입력 집계", "인터럽트 큐에서 흘린 수 + 버퍼에서 만료된 수를 합산 보고", AMBER),
    ("힙 최저치 추적", "가장 부족했던 순간의 여유 메모리를 기록 — 누수 조기 발견", MUTED),
]
y = Inches(2.25)
for name, desc, col in items:
    card(s, Inches(0.9), y, Inches(11.5), Inches(1.02))
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.3), y + Inches(0.38), Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = col
    dot.line.fill.background(); dot.shadow.inherit = False
    one(s, Inches(1.7), y + Inches(0.19), Inches(3.8), Inches(0.4), name, 19, True, INK)
    one(s, Inches(5.6), y + Inches(0.24), Inches(6.5), Inches(0.4), desc, 15, False, MUTED)
    y += Inches(1.18)
one(s, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.35),
    "숨기지 않는 것이 핵심입니다. 조용히 재시작하면 사용자는 자기 손을 의심하게 됩니다.",
    15, True, COBALT)

# --- 10. 오프라인 버퍼 -------------------------------------------
s = slide(prs)
header(s, "가장 나쁜 고장을 막다", "폰이 없을 때 누른 입력을 잃지 않습니다")
card(s, Inches(0.9), Inches(2.2), Inches(5.55), Inches(2.15))
one(s, Inches(1.3), Inches(2.48), Inches(4.8), Inches(0.4), "예전", 18, True, RED)
one(s, Inches(1.3), Inches(2.95), Inches(4.8), Inches(1.2),
    "링크가 끊겨도 그대로 내보내고\nACK 를 3회 기다리다 버림.\n그동안 다음 입력도 함께 사라짐.",
    15, False, MUTED, SANS, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 1.5)
card(s, Inches(6.85), Inches(2.2), Inches(5.55), Inches(2.15))
one(s, Inches(7.25), Inches(2.48), Inches(4.8), Inches(0.4), "지금", 18, True, GREEN)
one(s, Inches(7.25), Inches(2.95), Inches(4.8), Inches(1.2),
    "16칸 고리 버퍼에 담아 두고\n재연결되면 누른 순서대로 되살림.\n밀려난 수는 통계로 보고.",
    15, False, MUTED, SANS, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 1.5)
card(s, Inches(0.9), Inches(4.65), Inches(11.5), Inches(1.75), WHITE, AMBER)
one(s, Inches(1.35), Inches(4.92), Inches(10.6), Inches(0.4),
    "설계 판단 — 무작정 되살리지 않습니다", 19, True, AMBER)
one(s, Inches(1.35), Inches(5.4), Inches(10.6), Inches(0.9),
    "10분 전에 누른 '전화'가 재연결되는 순간 갑자기 걸리면 그것은 고장보다 나쁩니다.\n"
    "10초 안쪽만 되살리고, 그보다 오래된 것은 실행하지 않고 '놓친 입력' 수로만 보고합니다.",
    15, False, MUTED, SANS, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 1.5)

# --- 11. 저전력 --------------------------------------------------
s = slide(prs)
header(s, "저전력", "하루 종일 꽂아 두는 물건이므로")
one(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.4),
    "USB 기기가 잠드는 것은 위험합니다. 그래서 조건을 좁혔습니다.", 17, True, INK)
card(s, Inches(0.9), Inches(2.6), Inches(5.55), Inches(2.5))
one(s, Inches(1.3), Inches(2.88), Inches(4.8), Inches(0.4), "잠드는 조건 (모두 만족)", 17, True, COBALT)
one(s, Inches(1.3), Inches(3.35), Inches(4.8), Inches(1.6),
    "· 폰이 듣고 있지 않음 (DTR 내려감)\n· 버튼이 눌려 있지 않음\n· 보낼 것이 남아 있지 않음\n· 3초간 조용함",
    15, False, MUTED, SANS, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 1.5)
card(s, Inches(6.85), Inches(2.6), Inches(5.55), Inches(2.5))
one(s, Inches(7.25), Inches(2.88), Inches(4.8), Inches(0.4), "깨어나는 조건", 17, True, GREEN)
one(s, Inches(7.25), Inches(3.35), Inches(4.8), Inches(1.6),
    "· 버튼 눌림 (GPIO 로우 레벨)\n· 타이머 0.5초\n  → 깨어나 DTR 을 다시 확인\n  → 앱 재연결을 0.5초 안에 인지",
    15, False, MUTED, SANS, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 1.5)
one(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.8),
    "한 번에 0.5초씩만 자는 이유는 워치독입니다. 자는 동안에는 어느 태스크도 워치독을 먹이지\n"
    "못하므로, 슬라이스가 워치독 시간(5초)을 넘으면 멀쩡한 기기가 재시작해 버립니다.",
    15, False, MUTED, SANS, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 1.5)
one(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.4),
    f"대기 전류  {fmt(MEASURED['idle_current_ma'],'mA')}  →  절전 {fmt(MEASURED['sleep_current_ua'],'µA',0)}",
    18, True, COBALT if MEASURED["idle_current_ma"] else LIGHT, MONO)

# --- 12. 검증 ----------------------------------------------------
s = slide(prs)
header(s, "검증", "지어낸 숫자는 없습니다")
stats = [
    (f"{FACTS['tests_total']}", "개", "자동 테스트", f"펌웨어 {FACTS['tests_fw']} + 앱 {FACTS['tests_app']}"),
    (f"{FACTS['flash_pct']:.1f}", "%", "플래시 사용", f"{FACTS['flash_used']:,} / {FACTS['flash_total']:,} B"),
    (f"{FACTS['ram_pct']:.1f}", "%", "RAM 사용", f"{FACTS['ram_used']:,} / {FACTS['ram_total']:,} B"),
    (f"{FACTS['bom_won']:,}", "원", "부품값", f"{FACTS['bom_items']}종 · XIAO ESP32-S3 포함"),
]
x = Inches(0.9)
for val, unit, label, sub in stats:
    card(s, x, Inches(2.15), Inches(2.72), Inches(1.85))
    text(s, x, Inches(2.42), Inches(2.72), Inches(0.92),
         [[(val, 40, True, COBALT, SANS), (unit, 20, True, COBALT, SANS)]], PP_ALIGN.CENTER)
    one(s, x, Inches(3.18), Inches(2.72), Inches(0.35), label, 16, True, INK, SANS, PP_ALIGN.CENTER)
    one(s, x, Inches(3.55), Inches(2.72), Inches(0.35), sub, 10.5, False, MUTED, MONO, PP_ALIGN.CENTER)
    x += Inches(2.93)
rule(s, Inches(4.35))
one(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(0.4),
    "실기기 측정 (HIL)", 18, True, INK)
hil = [
    ("접점 → USB 송출 지연", fmt(MEASURED["latency_avg_ms"], "ms"), "목표 < 5 ms"),
    ("프레임 오류율", fmt(MEASURED["frame_error_pct"], "%"), "목표 < 0.1 %"),
    ("오실행", "0 건" if MEASURED["misfire_count"] == 0 else fmt(MEASURED["misfire_count"], "건", 0), "목표 0 건"),
]
y = Inches(5.05)
for name, val, goal in hil:
    done = "측정 예정" not in val
    one(s, Inches(0.95), y, Inches(4.6), Inches(0.35), name, 15, False, INK)
    one(s, Inches(5.7), y, Inches(3.0), Inches(0.35), val, 15, True, COBALT if done else LIGHT, MONO)
    one(s, Inches(9.0), y, Inches(3.3), Inches(0.35), goal, 13, False, LIGHT)
    y += Inches(0.5)
one(s, Inches(0.9), Inches(6.68), Inches(11.5), Inches(0.35),
    "측정 도구: firmware/tools/measure.py — 앱과 같은 규칙으로 CRC 검사·ACK 응답",
    13, False, LIGHT, MONO)

# --- 13. 2.0 → 3.1 -----------------------------------------------
s = slide(prs)
header(s, "규정 제10조 ③ 개선 사항", "2.0 대비 무엇이 달라졌나")
diffs = [
    ("기기 역할", "평문 문자열을 흘려보내는 전선", "프레임을 조립·서명·재전송하는 노드"),
    ("입력 처리", "폴링 + 손 뗄 때까지 블로킹 대기", "인터럽트 + FreeRTOS 태스크 2개"),
    ("짧게/길게", "구조적으로 구분 불가", "0.7초 기준으로 구분 (버튼 3개 → 6가지)"),
    ("오류 검출", "없음", "CRC16 + ACK 재전송 + SEQ 중복 제거"),
    ("고장 대응", "없음", "워치독 5초 · 재시작 원인 보고"),
    ("입력 유실", "폰 없으면 사라짐", "16칸 버퍼 + 조건부 재생"),
    ("전력", "항상 활성", "조건부 light sleep"),
    ("검증", "없음", f"자동 테스트 {FACTS['tests_total']}개"),
]
y = Inches(1.95)
one(s, Inches(0.95), y, Inches(2.4), Inches(0.3), "항목", 13, True, COBALT)
one(s, Inches(3.5), y, Inches(4.0), Inches(0.3), "2.0", 13, True, LIGHT)
one(s, Inches(7.8), y, Inches(4.5), Inches(0.3), "3.1", 13, True, COBALT)
y += Inches(0.38)
for k, before, after in diffs:
    one(s, Inches(0.95), y, Inches(2.4), Inches(0.35), k, 14, True, INK)
    one(s, Inches(3.5), y, Inches(4.1), Inches(0.35), before, 13, False, LIGHT)
    one(s, Inches(7.8), y, Inches(4.5), Inches(0.35), after, 13, False, MUTED)
    y += Inches(0.6)

# --- 14. 마무리 --------------------------------------------------
s = slide(prs)
header(s, "정리", "왜 이것이 임베디드 소프트웨어인가")
pts = [
    ("입력을 만드는 쪽이 기기입니다", "화면이 아니라 물리 접점에서 시작하고, 그 순간을 인터럽트로 잡습니다."),
    ("통신 규격을 직접 정했습니다", "프레임·검사·재전송·중복 제거를 양쪽에 같은 규칙으로 구현했습니다."),
    ("기기가 스스로를 책임집니다", "멈추면 되살아나고, 왜 멈췄는지 보고하고, 못 보낸 입력을 쥐고 있습니다."),
    ("쓰는 사람이 어르신입니다", "기술 선택마다 '틀리게 동작하느니 아무것도 안 하는 편이 낫다'를 택했습니다."),
]
y = Inches(2.1)
for head, sub in pts:
    one(s, Inches(0.95), y, Inches(0.4), Inches(0.4), "—", 20, True, COBALT)
    one(s, Inches(1.5), y, Inches(10.8), Inches(0.4), head, 21, True, INK)
    one(s, Inches(1.5), y + Inches(0.44), Inches(10.8), Inches(0.4), sub, 15, False, MUTED)
    y += Inches(1.12)
rule(s, Inches(6.55))
one(s, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.4),
    f"{FACTS['product']} · 팀 {FACTS['team']} · {FACTS['contest']}", 14, False, LIGHT)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "VibeKey_발표자료.pptx")
prs.save(out)
print(f"저장: {out}")
print(f"장표 {len(prs.slides.__iter__.__self__._sldIdLst)}장")
